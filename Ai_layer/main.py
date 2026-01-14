from fastapi import FastAPI, File, UploadFile, HTTPException
import shutil
import os
import cv2
import pytesseract
import numpy as np
from pdf2image import convert_from_path
import re
import ollama
import docx
import json
import pandas as pd
import pptx

app = FastAPI()

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)  # Create uploads folder if not exists

# PDF TO IMAGE
def pdf_to_images(pdf_path, dpi=300):
    images = convert_from_path(pdf_path, dpi=dpi)
    opencv_images = []
    for img in images:
        img_np = np.array(img)
        img_cv = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)
        opencv_images.append(img_cv)
    return opencv_images

# ANGLE FINDER
def get_skew_angle(image):
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    gray = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(gray, 50, 150, apertureSize=3)
    lines = cv2.HoughLinesP(edges, 1, np.pi / 180, 100, minLineLength=100, maxLineGap=10)
    angles = []
    if lines is not None:
        for line in lines:
            x1, y1, x2, y2 = line[0]
            angle = np.arctan2(y2 - y1, x2 - x1) * 180 / np.pi
            angles.append(angle)
    return np.median(angles) if angles else 0

# ROTATE IMAGE
def rotate_image(image, angle):
    (h, w) = image.shape[:2]
    center = (w // 2, h // 2)
    M = cv2.getRotationMatrix2D(center, angle, 1.0)
    return cv2.warpAffine(image, M, (w, h), flags=cv2.INTER_CUBIC, borderMode=cv2.BORDER_REPLICATE)

# FIX SKEW
def correct_skew(image):
    return rotate_image(image, get_skew_angle(image))

# REMOVE ORIGINAL BORDER
def remove_borders(image):
    contours, hierarchy = cv2.findContours(image, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    cntsSorted = sorted(contours, key=lambda x: cv2.contourArea(x))
    cnt = cntsSorted[-1]
    x, y, w, h = cv2.boundingRect(cnt)
    return image[y:y+h, x:x+w]

# GRAYSCALE
def grayscale(image):
    return cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)

# TEXT CLEANING
def clean_text(text):
    lines = text.split('\n')
    cleaned_lines = [re.sub(r'[^a-zA-Z0-9,.\-/: ]+', ' ', line).strip() for line in lines if line.strip()]
    return '\n'.join(cleaned_lines)

# METADATA GENERATION
def generate_metadata(text, file_name=""):
    prompt = f"""
    You are an AI assistant specializing in metadata extraction.  
    Your primary task is to **accurately determine the document type** before extracting details.  
    Always ensure the **document type is correct**, using **text patterns, keywords, and structure** for classification.  
    If uncertain, choose the **most probable** category and justify why.

    ---

    ### **Step 1: Identify the Document Type**  
    Use the following classification logic:  

    - **Invoice** → If it contains words like `"Invoice No."`, `"Total Amount Due"`, `"Bill To"`.  
    - **Contract** → If it has `"Agreement between"`, `"Terms & Conditions"`, `"This contract is made on"`.  
    - **Email** → If it starts with `"Dear [Name]"`, contains `"From: [Email]"`, `"Subject:"`.  
    - **Research Paper** → If it has `"Abstract"`, `"Methodology"`, `"References"`, structured sections.  
    - **News Article** → If it has `"By [Author]"`, `"Published on"`, `"Dateline"`, `"Headline"`.  
    - **Notes** → If it is **unstructured** with **bullets, free text, short phrases**.  
    - **Report** → If it contains `"Executive Summary"`, `"Findings"`, `"Conclusion"`, `"Prepared by"`.  
    - **Legal Document** → If it includes `"Hereby"`, `"Witnesseth"`, `"Legal Code Section"`.  
    - **Attendance Sheet** → If it has `"Names"`, `"Sign-in"`, `"Date"`, `"Present/Absent"` as a **table or list**.  
    - **Literature (Poetry, Story, Essay, etc.)** → If it contains **poetic structure, storytelling elements, paragraphs of creative writing**.  
    - **Other** → If none of the above applies, classify as `"Other"` and **explain why**.  

    ---

    ### **Step 2: Extract Metadata Fields**  
    1. **Document Type (Highly Important)**  
    - Identify the **exact document type** based on the rules above.  
    - Justify your classification with **keywords and structure found in the text**.  

    2. **Key Entities**  
    - **People**: Names of individuals mentioned.  
    - **Organizations**: Companies, institutions, or groups referenced.  
    - **Dates**: Any specific dates mentioned.  

    3. **Summary**  
    - Provide a concise summary in **two sentences or less**.  

    4. **Additional Important Information**  
    - Any crucial monetary values, legal obligations, deadlines, or notable insights.  

    ---

    ### **Text to analyze:**  
    \"\"\"{text}\"\"\"


    GIVE OUTPUT IN JSON FORMAT ONLY
    """
    response = ollama.chat(model="gemma:2b", messages=[{"role": "user", "content": prompt}])
    return response['message']['content']

# DOCUMENT PROCESSING FUNCTIONS
def docx_to_text(docx_path):
    return "\n".join([para.text for para in docx.Document(docx_path).paragraphs])

def txt_to_text(txt_path):
    with open(txt_path, "r", encoding="utf-8") as file:
        return file.read()

def json_to_text(json_path):
    with open(json_path, "r", encoding="utf-8") as file:
        data = json.load(file)
        return "\n".join([f"{key}: {value}" for key, value in data.items()])

def excel_to_text(excel_path):
    df = pd.read_excel(excel_path, engine="openpyxl")
    return "\n".join(df.to_csv(index=False, sep="\t").split("\n"))

def ppt_to_text(ppt_path):
    text = []
    for slide in pptx.Presentation(ppt_path).slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)


# OCR FROM IMAGE
def image_to_text(image):
    if image is None:
        return ""
    ima = correct_skew(image)
    gray_image = grayscale(ima)
    _, im_bw = cv2.threshold(gray_image, 210, 230, cv2.THRESH_BINARY)
    no_borders = remove_borders(im_bw)
    color = [255, 255, 255]
    top, bottom, left, right = [150] * 4
    image_with_border = cv2.copyMakeBorder(no_borders, top, bottom, left, right, cv2.BORDER_CONSTANT, value=color)
    return clean_text(pytesseract.image_to_string(image_with_border))

# FILE TEXT EXTRACTION
def extract_text_from_file(file_path):
    ext = file_path.split(".")[-1].lower()
    full_ocr_text = ""
    if ext == "pdf":
        images = pdf_to_images(file_path)
        for idx, image in enumerate(images):
            if image is None:
                continue
            full_ocr_text += f"\n===== Page {idx} =====\n{image_to_text(image)}\n"
    elif ext in ["jpg", "jpeg", "png"]:
        full_ocr_text = image_to_text(cv2.imread(file_path))
    elif ext in ["doc", "docx"]:
        full_ocr_text = docx_to_text(file_path)
    elif ext in ["txt", "md"]:
        full_ocr_text = txt_to_text(file_path)
    elif ext == "json":
        full_ocr_text = json_to_text(file_path)
    elif ext in ["xls", "xlsx", "gsheet"]:
        full_ocr_text = excel_to_text(file_path)
    elif ext in ["ppt", "pptx", "gslides"]:
        full_ocr_text = ppt_to_text(file_path)
    return clean_text(full_ocr_text)

@app.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    file_path = os.path.join(UPLOAD_DIR, file.filename)

    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    extracted_text = extract_text_from_file(file_path)
    metadata = generate_metadata(extracted_text)

    os.remove(file_path)  # Delete the file after processing

    
    return metadata

@app.get("/")
def read_root():
    return {"message": "FastAPI File Processor is running"}